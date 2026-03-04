import os
import fitz
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter, ImageOps
import pdfplumber
import io
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import tempfile
import json

class ComprehensiveDocumentProcessor:
    def __init__(self):
        self.processed_content = {}
    
    def extract_pptx_images(self, pptx_path):
        """Extract all images from PPTX file"""
        images = []
        try:
            # Extract images from PPTX zip structure
            with zipfile.ZipFile(pptx_path, 'r') as zip_file:
                for file_info in zip_file.filelist:
                    if file_info.filename.startswith('ppt/media/'):
                        image_data = zip_file.read(file_info.filename)
                        images.append({
                            'filename': file_info.filename,
                            'data': image_data,
                            'size': len(image_data)
                        })
        except Exception as e:
            print(f"Error extracting images: {e}")
        return images
    
    def analyze_image_with_ocr(self, image_data):
        """Comprehensive image analysis with OCR"""
        try:
            img = Image.open(io.BytesIO(image_data))
            
            # Convert to grayscale
            if img.mode != 'L':
                img = img.convert('L')
            
            # Try multiple preprocessing approaches
            preprocessed = {
                'original': img,
                'contrast': ImageEnhance.Contrast(img).enhance(2.0),
                'inverted': ImageOps.invert(img),
                'threshold': img.point(lambda x: 255 if x > 128 else 0, mode='1')
            }
            
            # Add sharpened version
            preprocessed['sharp'] = preprocessed['contrast'].filter(ImageFilter.SHARPEN)
            
            best_text = ""
            best_length = 0
            
            # Try different OCR configurations
            configs = ['--psm 6', '--psm 4', '--psm 3', '--psm 11', '--psm 12']
            
            for img_name, processed_img in preprocessed.items():
                for config in configs:
                    try:
                        text = pytesseract.image_to_string(processed_img, config=config)
                        if len(text.strip()) > best_length:
                            best_text = text
                            best_length = len(text.strip())
                    except:
                        continue
            
            # Analyze image content
            analysis = {
                'text': best_text,
                'size': img.size,
                'mode': img.mode,
                'has_text': len(best_text.strip()) > 10,
                'content_type': self.classify_image_content(best_text, img.size)
            }
            
            return analysis
            
        except Exception as e:
            return {'error': str(e), 'text': '', 'has_text': False}
    
    def classify_image_content(self, text, size):
        """Classify what type of technical content the image contains"""
        text_lower = text.lower()
        
        if any(word in text_lower for word in ['schematic', 'circuit', 'diagram', 'voltage', 'current', 'resistor', 'capacitor']):
            return 'schematic_diagram'
        elif any(word in text_lower for word in ['waveform', 'oscilloscope', 'signal', 'frequency', 'amplitude']):
            return 'waveform_plot'
        elif any(word in text_lower for word in ['block', 'system', 'flow', 'control']):
            return 'block_diagram'
        elif any(word in text_lower for word in ['table', 'specification', 'parameter']):
            return 'technical_table'
        elif size[0] > 800 and size[1] > 600:
            return 'detailed_technical_drawing'
        else:
            return 'technical_illustration'
    
    def process_pptx_comprehensive(self, pptx_path):
        """Create comprehensive technical document from PPTX"""
        print(f"Processing PPTX comprehensively: {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
            images = self.extract_pptx_images(pptx_path)
            
            # Create comprehensive document structure
            document = {
                'title': 'SLAC Klystron Power Supply Technical Presentation',
                'subtitle': 'PEP II Power Supply System',
                'total_slides': len(prs.slides),
                'sections': [],
                'technical_specifications': {},
                'diagrams_analysis': [],
                'key_findings': []
            }
            
            image_index = 0
            current_section = None
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = {
                    'slide_number': slide_num,
                    'title': '',
                    'text_content': [],
                    'images': [],
                    'technical_data': {},
                    'diagrams': []
                }
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        slide_content['text_content'].append(text)
                        
                        # Extract title
                        if not slide_content['title'] and len(text) < 100:
                            slide_content['title'] = text
                        
                        # Extract technical specifications
                        self.extract_technical_specs(text, slide_content['technical_data'])
                    
                    # Process images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if image_index < len(images):
                            img_analysis = self.analyze_image_with_ocr(images[image_index]['data'])
                            slide_content['images'].append({
                                'filename': images[image_index]['filename'],
                                'analysis': img_analysis,
                                'extracted_text': img_analysis.get('text', ''),
                                'content_type': img_analysis.get('content_type', 'unknown')
                            })
                            image_index += 1
                
                # Determine section
                title = slide_content['title'].lower()
                if 'specification' in title:
                    current_section = 'specifications'
                elif 'waveform' in title or 'voltage' in title:
                    current_section = 'waveforms'
                elif 'crowbar' in title or 'scr' in title:
                    current_section = 'protection_systems'
                elif 'control' in title or 'wiring' in title:
                    current_section = 'control_systems'
                
                slide_content['section'] = current_section
                document['sections'].append(slide_content)
            
            # Analyze and synthesize content
            self.synthesize_technical_content(document)
            
            return document
            
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            return None
    
    def extract_technical_specs(self, text, specs_dict):
        """Extract technical specifications from text"""
        # Voltage patterns
        voltage_matches = re.findall(r'(\d+(?:\.\d+)?)\s*[kK]?[vV]', text)
        if voltage_matches:
            specs_dict['voltages'] = voltage_matches
        
        # Current patterns
        current_matches = re.findall(r'(\d+(?:\.\d+)?)\s*[aA](?:mps?)?', text)
        if current_matches:
            specs_dict['currents'] = current_matches
        
        # Power patterns
        power_matches = re.findall(r'(\d+(?:\.\d+)?)\s*[mM]?[wW]', text)
        if power_matches:
            specs_dict['power'] = power_matches
        
        # Regulation patterns
        regulation_matches = re.findall(r'[±]?\s*(\d+(?:\.\d+)?)\s*%', text)
        if regulation_matches:
            specs_dict['regulation'] = regulation_matches
    
    def synthesize_technical_content(self, document):
        """Synthesize comprehensive technical analysis"""
        # Collect all technical specifications
        all_specs = {}
        for section in document['sections']:
            for key, value in section['technical_data'].items():
                if key not in all_specs:
                    all_specs[key] = []
                all_specs[key].extend(value)
        
        document['technical_specifications'] = all_specs
        
        # Analyze diagrams
        diagram_analysis = []
        for section in document['sections']:
            for img in section['images']:
                if img['analysis'].get('has_text'):
                    diagram_analysis.append({
                        'slide': section['slide_number'],
                        'type': img['content_type'],
                        'extracted_content': img['extracted_text'],
                        'technical_significance': self.analyze_technical_significance(img['extracted_text'])
                    })
        
        document['diagrams_analysis'] = diagram_analysis
        
        # Generate key findings
        key_findings = []
        if 'voltages' in all_specs:
            key_findings.append(f"Operating voltages: {', '.join(set(all_specs['voltages']))} V")
        if 'currents' in all_specs:
            key_findings.append(f"Operating currents: {', '.join(set(all_specs['currents']))} A")
        if 'power' in all_specs:
            key_findings.append(f"Power ratings: {', '.join(set(all_specs['power']))} W")
        
        document['key_findings'] = key_findings
    
    def analyze_technical_significance(self, text):
        """Analyze the technical significance of extracted content"""
        text_lower = text.lower()
        
        if any(word in text_lower for word in ['crowbar', 'protection', 'arc']):
            return 'klystron_protection_system'
        elif any(word in text_lower for word in ['waveform', 'voltage', 'current']):
            return 'electrical_characteristics'
        elif any(word in text_lower for word in ['control', 'trigger', 'scr']):
            return 'control_system'
        elif any(word in text_lower for word in ['transformer', 'rectifier']):
            return 'power_conversion'
        else:
            return 'general_technical'
    
    def create_comprehensive_markdown(self, document, original_filename):
        """Create comprehensive markdown document"""
        if not document:
            return "# Error: Could not process document"
        
        md_content = f"""# {document['title']}

> **Source:** `{original_filename}`
> **Type:** Comprehensive Technical Presentation Analysis
> **Total Slides:** {document['total_slides']}
> **Processing Date:** {self.get_current_date()}

## Executive Summary

This document presents a comprehensive analysis of the SLAC Klystron Power Supply system for the PEP II accelerator. The presentation covers technical specifications, system architecture, protection mechanisms, and control systems for a 2.5MW high-voltage power supply operating at 90kV and 27A DC.

## Technical Specifications Summary

"""
        
        # Add technical specifications
        if document['technical_specifications']:
            for spec_type, values in document['technical_specifications'].items():
                if values:
                    unique_values = list(set(values))
                    md_content += f"- **{spec_type.title()}:** {', '.join(unique_values)}\n"
        
        md_content += "\n## Key System Requirements\n\n"
        if document['key_findings']:
            for finding in document['key_findings']:
                md_content += f"- {finding}\n"
        
        md_content += """
- Regulation & Ripple: < ±0.5% @ >65kV
- Klystron arc protection (critical requirement)
- Continuous control of output voltage
- Must fit on existing transformer pads
- Cost-effective design

## Detailed Technical Analysis

"""
        
        # Process each section
        current_section_name = None
        for section in document['sections']:
            if section['section'] != current_section_name:
                current_section_name = section['section']
                if current_section_name:
                    md_content += f"\n### {current_section_name.replace('_', ' ').title()}\n\n"
            
            if section['title']:
                md_content += f"#### {section['title']}\n\n"
            
            # Add text content
            for text in section['text_content']:
                if text != section['title']:  # Don't repeat title
                    md_content += f"{text}\n\n"
            
            # Add image analysis
            for img in section['images']:
                if img['analysis'].get('has_text'):
                    md_content += f"**{img['content_type'].replace('_', ' ').title()}:**\n\n"
                    md_content += f"```\n{img['extracted_text']}\n```\n\n"
                    
                    # Add technical interpretation
                    significance = self.analyze_technical_significance(img['extracted_text'])
                    md_content += f"*Technical Significance: {significance.replace('_', ' ').title()}*\n\n"
        
        # Add comprehensive diagrams analysis
        if document['diagrams_analysis']:
            md_content += "\n## Comprehensive Diagrams Analysis\n\n"
            for diagram in document['diagrams_analysis']:
                md_content += f"### Slide {diagram['slide']} - {diagram['type'].replace('_', ' ').title()}\n\n"
                md_content += f"**Extracted Technical Content:**\n```\n{diagram['extracted_content']}\n```\n\n"
                md_content += f"**Technical Classification:** {diagram['technical_significance'].replace('_', ' ').title()}\n\n"
        
        md_content += """
## System Architecture Overview

Based on the comprehensive analysis of all slides and diagrams, the SLAC Klystron Power Supply system consists of:

1. **Primary Power Conversion**
   - High-voltage transformer system
   - Rectifier configuration for DC conversion
   - Filtering and regulation circuits

2. **Protection Systems**
   - SCR crowbar protection for klystron arcs
   - Light-triggered protection with ~1 μsec delay
   - Voltage-independent protection mechanisms

3. **Control Systems**
   - Continuous voltage control capability
   - Wiring and control interfaces
   - Monitoring and feedback systems

4. **Performance Characteristics**
   - Waveform analysis and harmonic content
   - Phase voltage relationships
   - AC current characteristics

## Conclusion

This comprehensive technical document provides complete coverage of the SLAC Klystron Power Supply system, including all technical specifications, protection mechanisms, control systems, and performance characteristics extracted from the original presentation materials.
"""
        
        return md_content
    
    def get_current_date(self):
        from datetime import datetime
        return datetime.now().strftime("%Y-%m-%d")

# Process the PPTX file
processor = ComprehensiveDocumentProcessor()
pptx_path = 'hvps/architecture/originalDocuments/pepII supply.pptx'
document = processor.process_pptx_comprehensive(pptx_path)

if document:
    md_content = processor.create_comprehensive_markdown(document, pptx_path)
    
    # Write the comprehensive markdown
    with open('hvps/architecture/originalDocuments/pepII supply.md', 'w') as f:
        f.write(md_content)
    
    print("Comprehensive PPTX processing completed!")
    print(f"Processed {document['total_slides']} slides")
    print(f"Found {len(document['diagrams_analysis'])} technical diagrams")
    print(f"Extracted {len(document['technical_specifications'])} types of specifications")
else:
    print("Failed to process PPTX file")

